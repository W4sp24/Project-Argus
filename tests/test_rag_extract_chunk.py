"""Tests for RAG extraction and chunking (no heavy deps required)."""

import importlib.util
from pathlib import Path

import pytest

from backend.rag.chunk import chunk_blocks
from backend.rag.extract import Block, extract_blocks
from backend.rag.paths import is_indexable

MINIMAL_PDF = b"""%PDF-1.4
1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj
2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj
3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
  /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj
4 0 obj << /Length 60 >> stream
BT /F1 12 Tf 72 720 Td (Dijkstra invented shortest paths) Tj ET
endstream endobj
5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj
trailer << /Root 1 0 R >>
%%EOF
"""


def test_is_indexable_enforces_privacy_and_types() -> None:
    assert is_indexable("10-Daily/2026-07-12.md")
    assert is_indexable("15-Courses/CS201/materials/deck.pptx")
    assert not is_indexable("99-Private/diary.md"), "I3 violation"
    assert not is_indexable("90-Meta/sessions/2026/x.md"), "journal is not RAG input"
    assert not is_indexable(".obsidian/workspace.md")
    assert not is_indexable("50-Reference/image.png")


def test_markdown_extraction_strips_frontmatter_and_skips_no_ai(tmp_path: Path) -> None:
    note = tmp_path / "note.md"
    note.write_text(
        "---\ntitle: Graphs\ntags: [cs, algs]\n---\n\n# Graphs\n\nBFS uses a queue.\n",
        encoding="utf-8",
    )
    blocks = extract_blocks(note)
    assert len(blocks) == 1
    assert "BFS uses a queue" in blocks[0].text
    assert "---" not in blocks[0].text
    assert blocks[0].meta["frontmatter"]["title"] == "Graphs"

    private = tmp_path / "private.md"
    private.write_text("Some text #no-ai here\n", encoding="utf-8")
    assert extract_blocks(private) == [], "I3 violation: no-ai note extracted"


@pytest.mark.skipif(importlib.util.find_spec("pdfplumber") is None, reason="rag extra missing")
def test_pdf_extraction_keeps_page_numbers(tmp_path: Path) -> None:
    pdf = tmp_path / "lecture.pdf"
    pdf.write_bytes(MINIMAL_PDF)
    blocks = extract_blocks(pdf)
    assert blocks, "no text extracted from PDF"
    assert blocks[0].meta["page"] == 1
    assert "Dijkstra" in blocks[0].text


@pytest.mark.skipif(importlib.util.find_spec("pptx") is None, reason="rag extra missing")
def test_pptx_extraction_keeps_slide_numbers(tmp_path: Path) -> None:
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Hash tables amortize to O(1)"
    path = tmp_path / "deck.pptx"
    deck.save(path)

    blocks = extract_blocks(path)
    assert blocks[0].meta["slide"] == 1
    assert "Hash tables" in blocks[0].text


def test_chunking_splits_headings_and_carries_metadata() -> None:
    text = "# Sorting\n\nIntro paragraph.\n\n## Quicksort\n\nPivot logic and [[Hoare]] scheme.\n"
    blocks = [Block(text=text, meta={"frontmatter": {"tags": ["cs"], "created": "2026-07-01"}})]

    chunks = chunk_blocks(blocks, "15-Courses/CS201/notes/sorting.md")

    headings = {chunk.meta["heading"] for chunk in chunks}
    assert "Quicksort" in headings
    sample = chunks[0].meta
    assert sample["course"] == "CS201"
    assert sample["tags"] == "cs"
    assert sample["date"] == "2026-07-01"
    assert any("Hoare" in chunk.meta["wikilinks"] for chunk in chunks)


def test_chunking_windows_long_sections_with_overlap() -> None:
    words = " ".join(f"word{i}" for i in range(600))
    chunks = chunk_blocks([Block(text=words, meta={"frontmatter": {}})], "note.md")

    assert len(chunks) >= 2
    first_words = chunks[0].text.split()
    second_words = chunks[1].text.split()
    assert first_words[-1] in second_words, "windows must overlap"


def test_chunking_passes_page_meta_through() -> None:
    chunks = chunk_blocks([Block(text="Slide text", meta={"page": 7})], "docs/x.pdf")
    assert chunks[0].meta["page"] == 7
