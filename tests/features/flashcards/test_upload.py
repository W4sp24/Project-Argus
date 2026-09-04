"""One-shot generation from a file the vault never sees.

The interesting case here is not the happy path -- it is the oversized block.
``.docx`` extraction returns a single block for the whole document, and
``deck_prompt`` tests each *whole* block against its 60k budget and breaks, so
one big block yields a prompt with an empty ``SOURCES:`` section. The corpus is
not empty, so ``generate_cards``' own guard does not fire either, and the model
writes plausible cards from nothing at all. That failure is silent, which is why
it has a test of its own.
"""

from __future__ import annotations

import importlib.util
import io
from pathlib import Path
from typing import Any

import pytest
from fastapi.testclient import TestClient

from backend.core.config import Settings
from backend.features.flashcards import generate, upload
from backend.features.flashcards.store import FlashcardsError
from backend.main import create_app

DOCX = pytest.mark.skipif(importlib.util.find_spec("docx") is None, reason="rag extra missing")
PPTX = pytest.mark.skipif(importlib.util.find_spec("pptx") is None, reason="rag extra missing")

REPLY = """\
Q:: what is a flow network
A:: a directed graph whose edges carry capacities
"""


def _generator(reply: str) -> Any:
    async def run(prompt: str) -> str:
        run.prompt = prompt  # type: ignore[attr-defined]
        return reply

    return run


def _docx_bytes(paragraphs: list[str]) -> bytes:
    from docx import Document

    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


# --- splitting -------------------------------------------------------------


def test_split_prefers_a_paragraph_then_a_line_then_a_word() -> None:
    assert upload.split_text("alpha beta\n\ngamma delta epsilon", 14) == [
        "alpha beta",
        "gamma delta",
        "epsilon",
    ]


def test_split_never_yields_an_empty_excerpt() -> None:
    assert upload.split_text("   \n\n   ", 8) == []
    assert all(piece.strip() for piece in upload.split_text("a b c   d", 3))


def test_split_breaks_an_unbroken_run_rather_than_looping_forever() -> None:
    assert upload.split_text("x" * 25, 10) == ["x" * 10, "x" * 10, "x" * 5]


# --- what the corpus looks like --------------------------------------------


def test_an_unsupported_suffix_is_refused_before_a_byte_is_read() -> None:
    """A .txt returns [] from extract_blocks *and* appends no error, so it is
    indistinguishable from a blank file. Those are read in the browser."""
    with pytest.raises(FlashcardsError) as caught:
        upload.corpus_from_upload("notes.txt", io.BytesIO(b"front\tback"))
    assert "not a kind Argus can read here" in str(caught.value)


def test_a_body_over_the_cap_is_refused(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(upload, "MAX_UPLOAD_BYTES", 16)
    monkeypatch.setattr(upload, "COPY_CHUNK_BYTES", 8)
    with pytest.raises(upload.UploadTooLargeError):
        upload.corpus_from_upload("big.pdf", io.BytesIO(b"z" * 64))


def test_an_empty_file_says_it_is_empty() -> None:
    with pytest.raises(FlashcardsError) as caught:
        upload.corpus_from_upload("blank.pdf", io.BytesIO(b""))
    assert "empty" in str(caught.value)


@DOCX
def test_a_docx_names_itself_as_the_source() -> None:
    corpus = upload.corpus_from_upload(
        "lecture-04.docx", io.BytesIO(_docx_bytes(["flow networks carry capacities"]))
    )
    assert [chunk["meta"]["path"] for chunk in corpus] == ["lecture-04.docx"]


@DOCX
def test_one_oversized_block_still_reaches_the_prompt(monkeypatch: pytest.MonkeyPatch) -> None:
    """The regression this module exists for.

    Unsplit, a document longer than ``MAX_PROMPT_CHARS`` fails ``deck_prompt``'s
    budget test on its first and only block, and the prompt ships with nothing
    under ``SOURCES:`` -- an ungrounded deck, generated in silence.
    """
    monkeypatch.setattr(upload, "MAX_BLOCK_CHARS", 2_000)
    long_paragraph = "capacities and cuts. " * 6_000  # ~126k chars, > MAX_PROMPT_CHARS
    corpus = upload.corpus_from_upload("long.docx", io.BytesIO(_docx_bytes([long_paragraph])))

    assert len(corpus) > 1
    assert all(len(chunk["text"]) <= 2_000 for chunk in corpus)
    sources = generate.deck_prompt("long", corpus, 10).split("SOURCES:", 1)[1]
    assert "[SOURCE path=long.docx" in sources
    assert "capacities and cuts" in sources


@PPTX
def test_slide_numbers_survive_into_the_corpus(tmp_path: Path) -> None:
    from pptx import Presentation

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "residual graphs"
    path = tmp_path / "week3.pptx"
    deck.save(path)

    corpus = upload.corpus_from_upload("week3.pptx", io.BytesIO(path.read_bytes()))
    assert corpus[0]["meta"]["slide"] == 1
    assert "slide 1" in generate.deck_prompt("week3", corpus, 5)


# --- the route -------------------------------------------------------------


@pytest.fixture()
def client(tmp_path: Path) -> TestClient:
    vault = tmp_path / "vault"
    (vault / "15-Courses" / "CS201" / "materials").mkdir(parents=True)
    return TestClient(
        create_app(
            Settings(_vault_path=vault),
            generator=_generator(REPLY),
            # Deterministic: run the job inline rather than on a daemon thread.
            ingest_job_runner=lambda run: run(),
        )
    )


@DOCX
def test_an_uploaded_file_becomes_a_deck_that_belongs_to_no_course(client: TestClient) -> None:
    response = client.post(
        "/api/flashcards/decks/generate/upload",
        files={"file": ("flow-networks.docx", _docx_bytes(["capacities and cuts"]))},
        data={"styles": ["definition", "concept"], "difficulty": "hard", "n": "5"},
    )
    assert response.status_code == 202, response.text

    deck = client.get(f"/api/flashcards/decks/{response.json()['deck_id']}").json()
    assert deck["title"] == "flow-networks"
    assert deck["course"] == ""
    assert deck["source"] == "generated"
    # Its own name, not a vault path -- which is exactly why it cannot collide
    # with a real source and wrongly badge a row in the SOURCES rail.
    assert deck["source_paths"] == ["flow-networks.docx"]
    assert deck["cards"] == 1
    assert deck["description"] == "hard · definition, concept · up to 5 cards"


@DOCX
def test_an_uploaded_file_can_still_be_filed_under_a_course(client: TestClient) -> None:
    response = client.post(
        "/api/flashcards/decks/generate/upload",
        files={"file": ("wk1.docx", _docx_bytes(["capacities"]))},
        data={"course": "CS201", "title": "week one"},
    )
    assert response.status_code == 202, response.text
    deck = client.get(f"/api/flashcards/decks/{response.json()['deck_id']}").json()
    assert deck["course"] == "CS201"
    assert deck["title"] == "week one"


def test_the_route_refuses_a_file_it_cannot_read(client: TestClient) -> None:
    response = client.post(
        "/api/flashcards/decks/generate/upload",
        files={"file": ("cards.txt", b"front\tback")},
    )
    assert response.status_code == 422
    assert "not a kind Argus can read here" in response.json()["detail"]


def test_the_route_refuses_an_unknown_difficulty_before_reading_the_file(
    client: TestClient,
) -> None:
    response = client.post(
        "/api/flashcards/decks/generate/upload",
        files={"file": ("x.docx", b"not really a docx")},
        data={"difficulty": "brutal"},
    )
    assert response.status_code == 422
    assert "unknown difficulty" in response.json()["detail"]


def test_a_file_that_cannot_be_parsed_says_so_rather_than_reporting_no_text(
    client: TestClient,
) -> None:
    """extract_blocks turns every failure into [] plus a log line, so without
    surfacing `errors` a corrupt file reads as an empty one."""
    response = client.post(
        "/api/flashcards/decks/generate/upload",
        files={"file": ("broken.docx", b"not really a docx")},
    )
    assert response.status_code == 422
    assert "could not read broken.docx" in response.json()["detail"]


def test_options_serve_the_upload_limits_rather_than_the_ui_mirroring_them(
    client: TestClient,
) -> None:
    options = client.get("/api/flashcards/generate/options").json()
    assert options["upload_suffixes"] == list(upload.UPLOAD_SUFFIXES)
    assert options["max_upload_bytes"] == upload.MAX_UPLOAD_BYTES
