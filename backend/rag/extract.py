"""Turn vault files into plain-text blocks with citation metadata.

Non-markdown course materials keep their page/slide numbers so downstream
answers can cite "file p.N" / "slide N" (invariant I6).
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import frontmatter

from backend.vault.privacy import is_no_ai

logger = logging.getLogger("argus.rag")


@dataclass
class Block:
    """One extractable unit of text plus citation metadata."""

    text: str
    meta: dict[str, Any] = field(default_factory=dict)


def _extract_markdown(file_path: Path) -> list[Block]:
    # No try/except here: extract_blocks already wraps every extractor call in
    # one, so failures are logged and collected in exactly one place.
    post = frontmatter.load(file_path)
    if is_no_ai(post):
        return []  # I3: tagged notes never enter the pipeline
    if not post.content.strip():
        return []
    return [Block(text=post.content, meta={"frontmatter": dict(post.metadata)})]


def _extract_pdf(file_path: Path) -> list[Block]:
    import pdfplumber  # heavy import kept lazy

    blocks: list[Block] = []
    with pdfplumber.open(file_path) as pdf:
        for number, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                blocks.append(Block(text=text, meta={"page": number}))
    return blocks


def _extract_pptx(file_path: Path) -> list[Block]:
    from pptx import Presentation  # heavy import kept lazy

    blocks: list[Block] = []
    for number, slide in enumerate(Presentation(file_path).slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if texts:
            blocks.append(Block(text="\n".join(texts), meta={"slide": number}))
    return blocks


def _extract_docx(file_path: Path) -> list[Block]:
    import docx  # heavy import kept lazy

    paragraphs = [p.text for p in docx.Document(file_path).paragraphs if p.text.strip()]
    if not paragraphs:
        return []
    return [Block(text="\n".join(paragraphs), meta={})]


def _extract_eml(file_path: Path) -> list[Block]:
    """One block per email: headers worth searching, then the plain body.

    The ingest route has accepted ``.eml`` since it was written, but there was
    no extractor for it -- so a dropped email was saved to the vault and then
    indexed to zero chunks, surfacing as "saved -- indexing unavailable". The
    subject and sender are prepended to the block text because searching for
    who a mail was from is at least as common as searching its body.
    """
    from backend.rag.email import parse_email

    parsed = parse_email(file_path.read_text(encoding="utf-8", errors="replace"))
    body = str(parsed.get("body") or "").strip()
    header_lines = [
        f"{label}: {parsed[key]}"
        for label, key in (("Subject", "subject"), ("From", "sender"), ("Date", "date"))
        if parsed.get(key)
    ]
    text = "\n".join([*header_lines, "", body]).strip() if header_lines else body
    if not text:
        return []
    meta = {key: parsed[key] for key in ("subject", "sender", "date") if parsed.get(key)}
    return [Block(text=text, meta=meta)]


_EXTRACTORS = {
    ".md": _extract_markdown,
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
    ".docx": _extract_docx,
    ".eml": _extract_eml,
}


def extract_blocks(file_path: Path, *, errors: list[str] | None = None) -> list[Block]:
    """Extract text blocks from a supported file; unsupported types yield [].

    ``errors``, when given, receives one message on failure. The empty-list
    return is unchanged either way — a single unreadable file must never break
    a caller that extracts many (e.g. a full reindex) — but a caller that
    wants to know *why* a file came back empty (rather than assume it was
    legitimately blank) can pass a list and inspect it afterward.
    """
    extractor = _EXTRACTORS.get(file_path.suffix.lower())
    if extractor is None:
        return []
    try:
        return extractor(file_path)
    except Exception as exc:
        logger.warning("failed to extract %s: %s", file_path, exc)
        if errors is not None:
            errors.append(str(exc))
        return []
